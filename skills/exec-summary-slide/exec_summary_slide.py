"""
exec_summary_slide.py — Generate the monthly Executive Summary ROB slide.

Usage:
  python exec_summary_slide.py --config <config.json> --output <output.pptx>

Config JSON structure:
{
  "month_year": "April 2026",
  "fiscal_year": "FY26",
  "sections": [
    {
      "title": "Highlights",
      "items": [
        {"text": "Plain bullet point"},
        {"text": "Sub-bullet under previous", "level": 2},
        {"label": "Discovery refresh", "text": "workshop hosted by Lester/Eric w/ 44 attendees"},
        {"label": "Channel Mode", "text": "Prototype engagement getting good feedback"}
      ]
    },
    {
      "title": "Headwinds",
      "items": [
        {"text": "CAPE deeply engaged in CY26 H1 planning for MCS"}
      ]
    }
  ],
  "engagement": {
    "teams": [
      {
        "name": "Depth",
        "envision": ["BASF", "Hilton", "Lloyds"],
        "build": ["Barclays", "Blackrock"],
        "deploy": ["BASF", "Chanel"],
        "adopt": ["Adecco", "BASF"],
        "total": 40
      },
      {
        "name": "FDE",
        "envision": ["Citgo", "Moderna"],
        "build": ["Beckton Dickinson", "UHG"],
        "deploy": [],
        "adopt": ["Amica", "United Health"],
        "total": 40
      }
    ],
    "totals": {
      "envision": 16,
      "build": 32,
      "deploy": 32,
      "adopt": null,
      "grand_total": 80
    }
  }
}
"""
import sys
import json
import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
except ImportError:
    print("python-pptx not installed. Run: python -m pip install python-pptx")
    sys.exit(1)

# Color palette
TITLE_BLUE = RGBColor(68, 114, 196)
SECTION_BLUE = RGBColor(47, 84, 150)
BODY_BLACK = RGBColor(51, 51, 51)
TABLE_HEADER_BG = RGBColor(214, 228, 240)
TABLE_BORDER = RGBColor(143, 170, 220)
GRADIENT_PURPLE = RGBColor(123, 104, 174)
GRADIENT_BLUE = RGBColor(91, 155, 213)
FOOTER_GRAY = RGBColor(128, 128, 128)
WHITE = RGBColor(255, 255, 255)
LIGHT_PURPLE = RGBColor(200, 185, 230)


def add_shape(slide, left, top, width, height, color, prs=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_gradient_border(slide, left, width, height, color1, color2):
    """Simulate gradient with stacked thin rectangles."""
    steps = 20
    step_h = height // steps
    for i in range(steps):
        r = int(color1[0] + (color2[0] - color1[0]) * i / steps)
        g = int(color1[1] + (color2[1] - color1[1]) * i / steps)
        b = int(color1[2] + (color2[2] - color1[2]) * i / steps)
        color = RGBColor(max(0, min(255, r)), max(0, min(255, g)), max(0, min(255, b)))
        shape = add_shape(slide, left, Emu(int(i * step_h)), width, Emu(int(step_h + 1)), color)
        sp = shape._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))


def set_paragraph(p, text, font_size=12, bold=False, color=BODY_BLACK,
                  underline=False, font_name='Calibri'):
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    run.font.underline = underline
    return run


def add_bullet(p):
    pPr = p._p.get_or_add_pPr()
    for existing in pPr.findall(qn('a:buChar')):
        pPr.remove(existing)
    for existing in pPr.findall(qn('a:buNone')):
        pPr.remove(existing)
    buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
    pPr.append(buChar)


def generate_slide(config, output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # -- Decorative gradient borders --
    border_width = Inches(0.4)
    # Left border
    add_gradient_border(slide, 0, border_width, prs.slide_height,
                        (180, 160, 220), (123, 104, 174))
    # Right border
    add_gradient_border(slide, Emu(int(prs.slide_width - border_width)),
                        border_width, prs.slide_height,
                        (91, 155, 213), (180, 160, 220))

    # ===== LEFT COLUMN =====
    left_x = 0.6
    left_w = 5.9
    y = 0.3

    # Title: "Executive Summary"
    tb = add_text_box(slide, left_x, y, left_w, 0.6)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_paragraph(p, "Executive Summary", font_size=32, bold=True, color=TITLE_BLUE)
    y += 0.7

    # Sections (flexible — Highlights, Headwinds, etc.)
    for section in config.get("sections", []):
        # Section heading
        tb = add_text_box(slide, left_x, y, left_w, 0.4)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        set_paragraph(p, section["title"], font_size=22, bold=True, color=SECTION_BLUE)
        y += 0.45

        # Bullet items
        tb = add_text_box(slide, left_x, y, left_w, 0.2 * len(section.get("items", [])) + 0.5)
        tf = tb.text_frame
        tf.word_wrap = True

        for i, item in enumerate(section.get("items", [])):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            level = item.get("level", 1)
            font_size = 12 if level == 1 else 11

            # Indentation
            indent = Inches(0.25 * (level - 1))
            p.paragraph_format.left_indent = indent
            p.space_after = Pt(3)

            add_bullet(p)

            if "label" in item:
                # Underlined lead-in label
                set_paragraph(p, item["label"] + ": ", font_size=font_size,
                              bold=True, underline=True, color=BODY_BLACK)
                run = p.add_run()
                run.text = item.get("text", "")
                run.font.size = Pt(font_size)
                run.font.color.rgb = BODY_BLACK
                run.font.name = 'Calibri'
            else:
                set_paragraph(p, item.get("text", ""), font_size=font_size, color=BODY_BLACK)

        # Calculate approximate height used
        item_count = len(section.get("items", []))
        y += 0.22 * item_count + 0.3

    # ===== RIGHT COLUMN =====
    right_x = 6.8
    right_w = 6.0
    ry = 0.3

    # Engagement Progress title
    fy = config.get("fiscal_year", "FY26")
    tb = add_text_box(slide, right_x, ry, right_w, 0.5)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_paragraph(p, f"{fy} Engagement Progress", font_size=22, bold=True, color=TITLE_BLUE)
    ry += 0.6

    # Build engagement table
    engagement = config.get("engagement", {})
    teams = engagement.get("teams", [])
    totals = engagement.get("totals", {})
    stages = ["envision", "build", "deploy", "adopt"]
    stage_labels = ["Envision", "Build", "Deploy", "Adopt"]

    # Header row + team rows + total row
    n_rows = 1 + len(teams) + 1
    n_cols = 6  # Team, Envision, Build, Deploy, Adopt, Total

    # Calculate table height based on max customers in any team
    max_customers = 1
    for team in teams:
        for stage in stages:
            count = len(team.get(stage, []))
            if count > max_customers:
                max_customers = count

    row_height = max(1.5, max_customers * 0.18 + 0.3)
    table_height = 0.4 + row_height * len(teams) + 0.4  # header + data + totals

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(right_x), Inches(ry),
        Inches(right_w), Inches(min(table_height, 6.0))
    )
    table = table_shape.table

    # Column widths
    col_widths = [0.8, 1.1, 1.1, 1.1, 1.1, 0.8]
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)

    # Header row
    headers = ["Team"] + stage_labels + ["Total"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = BODY_BLACK
            p.font.name = 'Calibri'

    # Team rows
    for t_idx, team in enumerate(teams):
        row_idx = t_idx + 1
        # Team name
        cell = table.cell(row_idx, 0)
        cell.text = team["name"]
        cell.vertical_anchor = MSO_ANCHOR.TOP
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = BODY_BLACK
            p.font.name = 'Calibri'

        # Stage cells with customer names
        for s_idx, stage in enumerate(stages):
            cell = table.cell(row_idx, s_idx + 1)
            cell.vertical_anchor = MSO_ANCHOR.TOP
            customers = team.get(stage, [])
            tf = cell.text_frame
            tf.word_wrap = True
            for c_idx, cust in enumerate(customers):
                p = tf.paragraphs[0] if c_idx == 0 else tf.add_paragraph()
                p.text = cust
                p.font.size = Pt(10)
                p.font.color.rgb = BODY_BLACK
                p.font.name = 'Calibri'
                p.space_after = Pt(1)

        # Total
        cell = table.cell(row_idx, 5)
        cell.vertical_anchor = MSO_ANCHOR.TOP
        cell.text = str(team.get("total", ""))
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = BODY_BLACK
            p.font.name = 'Calibri'

    # Totals row
    total_row = len(teams) + 1
    cell = table.cell(total_row, 0)
    cell.text = "WH Total"
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = BODY_BLACK
        p.font.name = 'Calibri'

    stage_totals = [totals.get(s) for s in stages]
    for i, val in enumerate(stage_totals):
        cell = table.cell(total_row, i + 1)
        cell.text = str(val) if val is not None else ""
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = BODY_BLACK
            p.font.name = 'Calibri'

    cell = table.cell(total_row, 5)
    cell.text = str(totals.get("grand_total", ""))
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BODY_BLACK
        p.font.name = 'Calibri'

    # ===== FOOTER =====
    tb = add_text_box(slide, 0.6, 7.1, 5, 0.3)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    set_paragraph(p, "Classified as Microsoft Confidential", font_size=8, color=FOOTER_GRAY)

    # Save
    prs.save(output_path)
    print(f"Saved: {output_path}")


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description="Generate Executive Summary slide")
    parser.add_argument("--config", required=True, help="Path to config JSON file")
    parser.add_argument("--output", required=True, help="Output .pptx path")
    args = parser.parse_args()

    config = json.loads(Path(args.config).read_text(encoding='utf-8'))
    generate_slide(config, args.output)
