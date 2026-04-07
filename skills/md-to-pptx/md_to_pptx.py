"""
md_to_pptx.py — Convert a markdown file to a formatted PowerPoint presentation.
Usage: python md_to_pptx.py <input.md> <output.pptx>
"""
import sys
import re
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
except ImportError:
    print("python-pptx not installed. Run: python -m pip install python-pptx")
    sys.exit(1)

# Color palette
DARK_BLUE = RGBColor(0, 51, 102)
ACCENT_BLUE = RGBColor(0, 120, 212)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 245, 245)
DARK_GRAY = RGBColor(51, 51, 51)
MED_GRAY = RGBColor(128, 128, 128)

MAX_BULLETS_PER_SLIDE = 6

class SlideBuilder:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def _add_shape(self, slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def _add_bg(self, slide, color):
        shape = self._add_shape(slide, 0, 0, self.prs.slide_width, self.prs.slide_height, color)
        sp = shape._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)

    def _add_text_box(self, slide, left, top, width, height, text,
                      font_size=18, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.alignment = align
        return tf

    def _add_bullet_list(self, slide, left, top, width, height, items,
                         font_size=16, color=DARK_GRAY):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            # Strip markdown bold
            clean = re.sub(r'\*\*(.+?)\*\*', r'\1', item)
            p.text = clean
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = 'Calibri'
            p.space_after = Pt(8)
            pPr = p._p.get_or_add_pPr()
            for existing in pPr.findall(qn('a:buChar')):
                pPr.remove(existing)
            for existing in pPr.findall(qn('a:buNone')):
                pPr.remove(existing)
            buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
            pPr.append(buChar)
        return tf

    def _add_slide_title(self, slide, text):
        self._add_shape(slide, Inches(0.8), Inches(0.5), Inches(0.15), Inches(0.6), ACCENT_BLUE)
        self._add_text_box(slide, 1.15, 0.45, 10, 0.8, text,
                          font_size=28, bold=True, color=DARK_BLUE)
        self._add_shape(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.02),
                       RGBColor(220, 220, 220))

    def add_title_slide(self, title, subtitle=None, meta=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, DARK_BLUE)
        self._add_shape(slide, 0, 0, Inches(0.25), self.prs.slide_height, ACCENT_BLUE)
        self._add_text_box(slide, 1.2, 1.8, 11, 1.2, title,
                          font_size=42, bold=True, color=WHITE)
        self._add_shape(slide, Inches(1.2), Inches(3.2), Inches(2.5), Inches(0.06), ACCENT_BLUE)
        if subtitle:
            self._add_text_box(slide, 1.2, 3.5, 11, 1.0, subtitle,
                              font_size=28, color=RGBColor(180, 200, 220))
        if meta:
            self._add_text_box(slide, 1.2, 5.0, 11, 1.0, meta,
                              font_size=14, color=RGBColor(160, 180, 200))
        return slide

    def add_content_slide(self, title, bullets=None, body_text=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_title(slide, title)
        y = 1.6
        if body_text:
            self._add_text_box(slide, 0.8, y, 11.5, 1.5, body_text,
                              font_size=18, color=DARK_GRAY)
            y += 1.5
        if bullets:
            # Split into multiple slides if too many bullets
            for chunk_start in range(0, len(bullets), MAX_BULLETS_PER_SLIDE):
                chunk = bullets[chunk_start:chunk_start + MAX_BULLETS_PER_SLIDE]
                if chunk_start == 0:
                    self._add_bullet_list(slide, 0.8, y, 11.5, 4.5, chunk, font_size=17)
                else:
                    slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
                    self._add_slide_title(slide, f"{title} (continued)")
                    self._add_bullet_list(slide, 0.8, 1.6, 11.5, 4.5, chunk, font_size=17)
        return slide

    def add_table_slide(self, title, headers, rows):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_title(slide, title)
        n_rows = len(rows) + 1
        n_cols = len(headers)
        col_width = 11.5 / n_cols
        table_shape = slide.shapes.add_table(
            n_rows, n_cols, Inches(0.8), Inches(1.6),
            Inches(11.5), Inches(0.5 * n_rows)
        )
        table = table_shape.table
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.font.name = 'Calibri'
        for r_idx, row in enumerate(rows):
            for c_idx, val in enumerate(row):
                if c_idx < n_cols:
                    cell = table.cell(r_idx + 1, c_idx)
                    cell.text = re.sub(r'\*\*(.+?)\*\*', r'\1', val.strip())
                    bg = LIGHT_GRAY if r_idx % 2 == 0 else WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = bg
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(12)
                        p.font.color.rgb = DARK_GRAY
                        p.font.name = 'Calibri'
        return slide

    def add_section_slide(self, title, subtitle=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, DARK_BLUE)
        self._add_shape(slide, Inches(0.8), Inches(2.8), Inches(1.5), Inches(0.06), ACCENT_BLUE)
        self._add_text_box(slide, 0.8, 1.4, 11, 1.5, title,
                          font_size=36, bold=True, color=WHITE)
        if subtitle:
            self._add_text_box(slide, 0.8, 3.1, 11, 1.5, subtitle,
                              font_size=18, color=RGBColor(180, 200, 220))
        return slide

    def save(self, path):
        self.prs.save(path)
        print(f"Saved: {path} ({len(self.prs.slides)} slides)")


def parse_markdown(md_text):
    """Parse markdown into a structured list of content blocks."""
    lines = md_text.split('\n')
    blocks = []
    current_heading = None
    current_level = 0
    current_bullets = []
    current_body = []
    current_table_headers = []
    current_table_rows = []
    in_table = False
    in_code = False
    meta_lines = []
    title = None

    def flush():
        nonlocal current_bullets, current_body, current_table_headers, current_table_rows, in_table
        if in_table and current_table_headers:
            blocks.append({
                'type': 'table',
                'heading': current_heading,
                'level': current_level,
                'headers': current_table_headers,
                'rows': current_table_rows,
            })
            current_table_headers = []
            current_table_rows = []
            in_table = False
        if current_bullets:
            blocks.append({
                'type': 'bullets',
                'heading': current_heading,
                'level': current_level,
                'items': current_bullets[:],
            })
            current_bullets = []
        if current_body:
            text = '\n'.join(current_body).strip()
            if text:
                blocks.append({
                    'type': 'body',
                    'heading': current_heading,
                    'level': current_level,
                    'text': text,
                })
            current_body = []

    for line in lines:
        if line.strip().startswith('```'):
            in_code = not in_code
            continue
        if in_code:
            current_body.append(line)
            continue

        # Table
        if line.strip().startswith('|'):
            cells = [c.strip() for c in line.strip().strip('|').split('|')]
            if not in_table:
                current_table_headers = cells
                in_table = True
            elif all(c.replace('-', '').replace(':', '').strip() == '' for c in cells):
                pass
            else:
                current_table_rows.append(cells)
            continue
        elif in_table:
            flush()

        # Heading
        heading_match = re.match(r'^(#{1,3})\s+(.+)$', line)
        if heading_match:
            flush()
            level = len(heading_match.group(1))
            text = heading_match.group(2).strip()
            if level == 1 and title is None:
                title = text
            current_heading = text
            current_level = level
            continue

        # Metadata
        meta_match = re.match(r'^\*\*(.+?):\*\*\s*(.+)$', line)
        if meta_match and title and not blocks:
            meta_lines.append(f"{meta_match.group(1)}: {meta_match.group(2)}")
            continue

        # Bullet
        bullet_match = re.match(r'^\s*[-*]\s+(.+)$', line)
        if bullet_match:
            current_bullets.append(bullet_match.group(1))
            continue

        # Blockquote
        if line.strip().startswith('>'):
            text = line.strip().lstrip('>').strip()
            if text:
                current_body.append(text)
            continue

        # Separator
        if line.strip() == '---':
            flush()
            continue

        # Regular text
        if line.strip():
            current_body.append(line.strip())

    flush()
    return title, meta_lines, blocks


def convert(input_path, output_path):
    md = Path(input_path).read_text(encoding='utf-8')
    title, meta_lines, blocks = parse_markdown(md)
    builder = SlideBuilder()

    # Title slide
    meta_str = '  |  '.join(meta_lines) if meta_lines else None
    subtitle_parts = title.split('\u2014') if title and '\u2014' in title else None
    if subtitle_parts and len(subtitle_parts) == 2:
        builder.add_title_slide(subtitle_parts[0].strip(), subtitle_parts[1].strip(), meta_str)
    else:
        builder.add_title_slide(title or 'Presentation', meta=meta_str)

    # Content slides
    last_h1 = None
    pending_body = None

    for block in blocks:
        heading = block.get('heading', '')
        level = block.get('level', 2)

        if block['type'] == 'table':
            builder.add_table_slide(heading or 'Data', block['headers'], block['rows'])

        elif block['type'] == 'bullets':
            builder.add_content_slide(heading or 'Details',
                                     bullets=block['items'],
                                     body_text=pending_body)
            pending_body = None

        elif block['type'] == 'body':
            if level == 1:
                last_h1 = heading
            pending_body = block['text']

    # Flush any remaining body text as a final slide
    if pending_body:
        builder.add_content_slide(last_h1 or 'Summary', body_text=pending_body)

    builder.save(output_path)


if __name__ == '__main__':
    if len(sys.argv) != 3:
        print("Usage: python md_to_pptx.py <input.md> <output.pptx>")
        sys.exit(1)
    convert(sys.argv[1], sys.argv[2])
